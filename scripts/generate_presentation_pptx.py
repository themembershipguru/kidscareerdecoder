"""Build BCA final presentation PPTX (and PDF if LibreOffice available)."""
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DIAG = ROOT / "docs" / "assets" / "diagrams"
SHOTS = ROOT / "docs" / "assets" / "screenshots"
LOGO = ROOT / "docs" / "assets" / "cu-logo-online.png"
DATE_STR = "22 May 2026"
DESKTOP_PPTX = Path("/Users/shubhlata/Desktop/KidsCareerDecoder_Presentation.pptx")
REPO_PPTX = ROOT / "docs" / "KidsCareerDecoder_Presentation.pptx"
DESKTOP_PDF = Path("/Users/shubhlata/Desktop/KidsCareerDecoder_Presentation.pdf")
REPO_PDF = ROOT / "docs" / "KidsCareerDecoder_Presentation.pdf"


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "KidsCareerDecoder"
    sub = slide.placeholders[1]
    sub.text = (
        "Web-Based Child Aptitude Assessment\n"
        "and Career Insight Platform\n\n"
        "Neel Kapoor | BCA Final Year, Semester VI\n"
        "Chandigarh University\n"
        f"Guide: Rohit Jha\n"
        f"Date: {DATE_STR}"
    )
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.4), Inches(0.35), width=Inches(3.2))


def add_bullet_slide(prs, title, bullets, image_path=None, image_right=True):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.4))
    footer.text_frame.text = f"KidsCareerDecoder | {DATE_STR}"
    footer.text_frame.paragraphs[0].font.size = Pt(10)

    if image_path and Path(image_path).exists():
        if image_right:
            slide.shapes.add_picture(
                str(image_path), Inches(5.0), Inches(1.35), width=Inches(4.5)
            )
            body.width = Inches(4.3)
        else:
            slide.shapes.add_picture(
                str(image_path), Inches(0.5), Inches(1.35), width=Inches(4.8)
            )
            body.left = Inches(5.4)
            body.width = Inches(4.2)


def add_image_slide(prs, title, image_paths, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)).text_frame.text = title
    slide.shapes[0].text_frame.paragraphs[0].font.size = Pt(28)
    slide.shapes[0].text_frame.paragraphs[0].font.bold = True

    paths = [Path(p) for p in image_paths if Path(p).exists()]
    if len(paths) == 1:
        slide.shapes.add_picture(str(paths[0]), Inches(0.8), Inches(1.0), width=Inches(8.5))
    elif len(paths) >= 2:
        slide.shapes.add_picture(str(paths[0]), Inches(0.5), Inches(1.0), width=Inches(4.3))
        slide.shapes.add_picture(str(paths[1]), Inches(5.0), Inches(1.0), width=Inches(4.3))

    if caption:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.6))
        box.text_frame.text = caption
        box.text_frame.paragraphs[0].font.size = Pt(12)

    foot = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.4))
    foot.text_frame.text = f"KidsCareerDecoder | {DATE_STR}"
    foot.text_frame.paragraphs[0].font.size = Pt(10)


def try_export_pdf(pptx_path, pdf_path):
    for cmd in (
        ["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(pdf_path.parent), str(pptx_path)],
        ["/Applications/LibreOffice.app/Contents/MacOS/soffice", "--headless", "--convert-to", "pdf", "--outdir", str(pdf_path.parent), str(pptx_path)],
    ):
        try:
            subprocess.run(cmd, check=True, capture_output=True, timeout=120)
            generated = pdf_path.parent / (pptx_path.stem + ".pdf")
            if generated.exists() and generated != pdf_path:
                generated.rename(pdf_path)
            if pdf_path.exists():
                return True
        except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired):
            continue
    return False


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_bullet_slide(
        prs,
        "Problem Statement",
        [
            "Parents lack structured, child-friendly insight into strengths",
            "Generic quizzes are not traceable or visible to guardians",
            "Need: secure family accounts, quizzes, AI careers with regional context",
        ],
    )

    add_bullet_slide(
        prs,
        "Objectives",
        [
            "Parent / child / admin roles with JWT security",
            "Published quizzes and adaptive session flow",
            "Six aptitude stripes + AI / fallback careers",
            "Parent dashboards, child results, admin console",
        ],
    )

    add_bullet_slide(
        prs,
        "Literature & Gap",
        [
            "Holland RIASEC vs six transparent aptitude stripes",
            "LMS quizzes lack parent analytics + LLM careers",
            "Gap closed: traceable scoring + admin + JSON-bound AI",
        ],
    )

    add_bullet_slide(
        prs,
        "System Architecture (DFD Level 0)",
        [
            "Parent: login, reports",
            "Child: quiz",
            "Admin: manage",
            "PostgreSQL (SQL) + OpenAI/Claude (HTTPS)",
        ],
        DIAG / "fig_4_01_dfd_level0.png",
    )

    add_bullet_slide(
        prs,
        "Database Design (ER Diagram)",
        [
            "users, quizzes, questions, question_options",
            "quiz_sessions (scores_json, metadata_json)",
            "quiz_answers, careers, app_settings",
        ],
        DIAG / "fig_4_04_er_diagram.png",
    )

    add_image_slide(
        prs,
        "Design — DFD Level 1 & Use Cases",
        [DIAG / "fig_4_02_dfd_level1.png", DIAG / "fig_4_03_use_case.png"],
        "Processes 1.0–6.0 and actor use cases (hand-drawn)",
    )

    add_bullet_slide(
        prs,
        "Sequence: Complete Session + AI",
        [
            "POST /session/:id/complete",
            "Tally + normalize scores → AI profile",
            "UPDATE quiz_sessions → 200 OK + careers",
        ],
        DIAG / "fig_4_05_sequence_complete.png",
    )

    add_bullet_slide(
        prs,
        "Implementation Highlights",
        [
            "bcrypt + JWT authentication",
            "sessionController: preprocess → getAptitudeProfile",
            "CareerResultCard: child vs parent variant",
            "Admin runtime AI provider selection",
        ],
    )

    add_image_slide(
        prs,
        "Parent Interface",
        [SHOTS / "fig_5_05_parent_add_child.png", SHOTS / "fig_5_08_parent_dashboard.png"],
        "Add child and dashboard (Fig 5.1, 5.2)",
    )

    add_image_slide(
        prs,
        "Child Interface",
        [SHOTS / "fig_5_07_child_quiz_select.png", SHOTS / "fig_7_02_child_results.png"],
        "Quiz selection and AI results (Fig 5.3, 7.2)",
    )

    add_image_slide(
        prs,
        "Admin Interface",
        [SHOTS / "fig_5_10_admin_overview.png", SHOTS / "fig_5_06_admin_insights.png"],
        "Overview and insights (Fig 5.5–5.6)",
    )

    add_bullet_slide(
        prs,
        "AI & Data Preprocessing",
        [
            "Tally answers → six percentage scores",
            "Age + country → Claude or OpenAI JSON careers",
            "Python notebook mirrors server math",
            "fallback_only when API unavailable",
        ],
    )

    add_bullet_slide(
        prs,
        "Testing Summary",
        [
            "27 documented test cases — Pass",
            "Auth, session, AI, admin, security scenarios",
            "CORS, JWT 401/403, reset token validation",
        ],
    )

    add_image_slide(
        prs,
        "Results",
        [SHOTS / "fig_7_01_parent_child_report.png", SHOTS / "fig_7_02_child_results.png"],
        "Parent report and child Logical Explorer results",
    )

    add_bullet_slide(
        prs,
        "Future Enhancements",
        [
            "Mobile app and HttpOnly cookie sessions",
            "Automated CI tests and Hindi UI",
            "IRT-based adaptive difficulty",
            "Cloud deployment documentation",
        ],
    )

    # Thank you slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You"
    slide.placeholders[1].text = (
        "Questions?\n\n"
        "Neel Kapoor | O23BCA110261@cuchd.in\n"
        "Guide: Rohit Jha\n"
        f"Date: {DATE_STR}"
    )

    for out in (DESKTOP_PPTX, REPO_PPTX):
        prs.save(str(out))
        print(out)

    if try_export_pdf(DESKTOP_PPTX, DESKTOP_PDF):
        REPO_PDF.write_bytes(DESKTOP_PDF.read_bytes())
        print(DESKTOP_PDF)
        print(REPO_PDF)
    else:
        print("PDF: export manually in PowerPoint (File → Export → PDF) or install LibreOffice")


if __name__ == "__main__":
    main()
